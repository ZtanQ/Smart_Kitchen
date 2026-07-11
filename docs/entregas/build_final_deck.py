"""Genera docs/entregas/Entrega6_Proyecto_Final_Defensa.pptx — el deck consolidado
para la sustentacion oral de la Entrega 6 (Deliverable 6). Los decks previos
(Entrega2, Entrega3, Entrega6_ComponenteAvanzado) cubren cada hito por separado;
este es el unico que narra el proyecto completo de punta a punta (motivacion ->
trabajo futuro) en una sola historia continua, siguiendo docs/defense_narrative.md.

Ejecutar: python docs/entregas/build_final_deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[2]
OUT = Path(__file__).resolve().parent / "Entrega6_Proyecto_Final_Defensa.pptx"
FIG = ROOT / "reports" / "figures"
SHOTS = ROOT / "reports" / "figures" / "dashboard_screenshots"

NAVY = RGBColor(0x1D, 0x2B, 0x4F)
BLUE = RGBColor(0x2A, 0x3B, 0x6C)
RED = RGBColor(0xD6, 0x45, 0x45)
AMBER = RGBColor(0xE8, 0xA8, 0x3A)
GREEN = RGBColor(0x5C, 0xB5, 0x5C)
MUTED = RGBColor(0x5B, 0x66, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x2A, 0x2F, 0x3A)
BG = RGBColor(0xF4, 0xF5, 0xF9)

SW, SH = Inches(13.333), Inches(7.5)  # 16:9


def add_bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_header(slide, eyebrow, title, color=NAVY):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12), Inches(1.05))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = eyebrow.upper()
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0xB8, 0xC2, 0xE0)
    p0.font.name = "Segoe UI"
    p1 = tf.add_paragraph()
    p1.text = title
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "Segoe UI"
    return tb


def add_bullets(slide, items, left=0.6, top=1.6, width=12.1, height=5.4, size=16,
                 color=INK, bullet_color=None, bold_first=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = f"{lead}  "
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = bullet_color or NAVY
            r1.font.name = "Segoe UI"
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = "Segoe UI"
        else:
            p.text = f"•  {item}"
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.name = "Segoe UI"
        p.space_after = Pt(10)
    return tb


def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.font.name = "Segoe UI"


def add_image_safe(slide, path, left, top, width=None, height=None):
    if path.exists():
        kwargs = {}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(str(path), left, top, **kwargs)
        return True
    return False


def stat_card(slide, left, top, w, h, label, value, sub, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF7, 0xF8, 0xFC)
    box.line.color.rgb = color
    box.line.width = Pt(2)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = label.upper()
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = MUTED
    p0.font.name = "Segoe UI"
    p1 = tf.add_paragraph()
    p1.text = value
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = color
    p1.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = INK
    p2.font.name = "Segoe UI"


prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]

# ---------------------------------------------------------------- Slide 1: Title
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.6))
tf = tb.text_frame
tf.word_wrap = True
p0 = tf.paragraphs[0]
p0.text = "SMART KITCHEN INTELLIGENCE (SKI)"
p0.font.size = Pt(40)
p0.font.bold = True
p0.font.color.rgb = WHITE
p0.font.name = "Segoe UI"
p1 = tf.add_paragraph()
p1.text = "Entrega 6 — Proyecto Final y Defensa"
p1.font.size = Pt(20)
p1.font.color.rgb = RGBColor(0xB8, 0xC2, 0xE0)
p1.font.name = "Segoe UI"
p1.space_before = Pt(14)
p2 = tf.add_paragraph()
p2.text = "¿Qué impacto tiene la ubicación física de un alimento en la pérdida económica y de salud del hogar?"
p2.font.size = Pt(15)
p2.font.italic = True
p2.font.color.rgb = RGBColor(0xDF, 0xE4, 0xF2)
p2.font.name = "Segoe UI"
p2.space_before = Pt(18)

tb2 = s.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(11.5), Inches(1.0))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Arroyo Gamarra, Favio Enrique · Melgar Puertas, José Guillermo · Reyna Alvarado, Gabriel Alonso"
p.font.size = Pt(13)
p.font.color.rgb = WHITE
p.font.name = "Segoe UI"
p2b = tf2.add_paragraph()
p2b.text = "1ACC0211 · Data Visualization · NRC 18519 · UPC 2026-I · Prof. Carlos Adrián Alarcón Delgado"
p2b.font.size = Pt(11)
p2b.font.color.rgb = RGBColor(0xB8, 0xC2, 0xE0)
p2b.font.name = "Segoe UI"

# ---------------------------------------------------------------- Slide 2: Agenda
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Agenda", "De la pregunta de negocio a las recomendaciones")
items = [
    "01 · Problema de negocio y pregunta de investigación",
    "02 · Dataset y usuario objetivo",
    "03 · Pipeline Python — metodología end-to-end",
    "04 · Limpieza de datos (reglas P1–P6)",
    "05 · Modelado predictivo y control de fuga de datos",
    "06 · Componente avanzado: PCA + t-SNE",
    "07 · Dashboard: KPIs, contexto de gasto, gráficos principales",
    "08 · Insights, conclusiones y recomendaciones",
    "09 · Limitaciones y trabajo futuro",
]
add_bullets(s, items, size=17)
add_footer(s, "Smart Kitchen Intelligence (SKI) — Entrega 6")

# ---------------------------------------------------------------- Slide 3: Business problem
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Motivación", "El desperdicio de alimentos es prevenible, pero invisible", RED)
add_bullets(s, [
    ("Problema:", "casi ningún hogar tiene visibilidad agregada de sus propios patrones de compra, consumo y desperdicio — las decisiones de compra semanal se toman a ciegas."),
    ("Pregunta de investigación:", "¿qué impacto tiene la ubicación física (Refrigerador / Despensa / Estante) en la pérdida económica y de salud del hogar?"),
    ("Usuario objetivo:", "responsable de compras y planificación de comidas de un hogar de 2–4 personas; usa el dashboard una vez por semana antes de la lista de compras."),
    ("Por qué importa:", "sirve como prototipo replicable — la misma lógica aplica a supermercados, bancos de alimentos o cadenas de retail que necesiten priorizar dónde intervenir."),
], top=1.7, size=17)
add_footer(s, "docs/proposal.md — pregunta analítica y usuario objetivo")

# ---------------------------------------------------------------- Slide 4: Dataset
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Dataset", "Catálogo real + movimientos simulados y calibrados", BLUE)
stat_card(s, Inches(0.6), Inches(1.6), Inches(3.9), Inches(1.5), "Eventos totales", "25,819", "10 hogares · 90 días", NAVY)
stat_card(s, Inches(4.7), Inches(1.6), Inches(3.9), Inches(1.5), "Productos en catálogo", "50", "OpenFoodFacts (real)", GREEN)
stat_card(s, Inches(8.8), Inches(1.6), Inches(3.9), Inches(1.5), "Ubicaciones físicas", "3", "Refrigerador · Despensa · Estante", AMBER)
add_bullets(s, [
    ("Catálogo:", "50 productos reales de OpenFoodFacts (API pública ODbL) — calorías, proteínas, carbohidratos, NutriScore."),
    ("Movimientos:", "simulación estocástica (src/simulation.py) calibrada con patrones reales de compra de Instacart Online Grocery y reglas de vida útil de la USDA FoodKeeper App — no son números al azar."),
    ("household_id:", "identificador sintético (0–9) de 10 hogares simulados independientes — variable de segmentación, no un dato personal ni generalizable."),
    ("Por qué simulado y no 100% real:", "conseguir datos reales de inventario doméstico a este nivel de detalle no es viable en el tiempo del curso; se declara la naturaleza sintética explícitamente."),
], top=3.4, size=15)
add_footer(s, "docs/data_dictionary.md · docs/source_inventory.md · docs/ethics_note.md")

# ---------------------------------------------------------------- Slide 5: Pipeline
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Metodología", "Pipeline Python reproducible, paso a paso", BLUE)
steps = [
    "simulation.py — genera movements_raw.csv",
    "ingestion.py — cataloga desde OpenFoodFacts",
    "anomaly_injection.py — inyecta anomalías realistas (ejercita la limpieza)",
    "preprocessing.py — limpieza P1–P6 → inventory_v1.csv (25,819 × 17)",
    "features.py + reduction.py — matriz de 61 features, PCA + t-SNE",
    "export_star_schema.py — modelo en estrella para Tableau",
    "build_dashboard_html.py / build_dashboard_twb.py — dashboards finales",
]
top = 1.7
for i, step in enumerate(steps):
    y = Inches(top + i * 0.72)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.4), Inches(0.4))
    dot.fill.solid()
    dot.fill.fore_color.rgb = NAVY
    dot.line.fill.background()
    dtf = dot.text_frame
    dtf.paragraphs[0].text = str(i + 1)
    dtf.paragraphs[0].font.size = Pt(14)
    dtf.paragraphs[0].font.bold = True
    dtf.paragraphs[0].font.color.rgb = WHITE
    dtf.paragraphs[0].alignment = PP_ALIGN.CENTER
    dtf.margin_top = Emu(0)
    dtf.margin_bottom = Emu(0)
    tb = s.shapes.add_textbox(Inches(1.2), Inches(top + i * 0.72 - 0.05), Inches(11.3), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = step
    p.font.size = Pt(15.5)
    p.font.color.rgb = INK
    p.font.name = "Segoe UI"
add_footer(s, "docs/runbook.md — secuencia completa y reproducible desde la raíz del repo")

# ---------------------------------------------------------------- Slide 6: Cleaning rules
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Limpieza de datos", "6 reglas verificables, cada una con efecto cuantificado", GREEN)
rows = [
    ("P1", "Outlier calórico", "Umbral físico 900 kcal/100g", "0 outliers en la fuente actual"),
    ("P2", "Homologación de categorías", "dept_id numérico → category_name + location derivada", "6 categorías legibles"),
    ("P3", "Integridad referencial", "product_id huérfano en movements", "0 huérfanos"),
    ("P4", "Tipos de dato explícitos", "cast antes de cualquier operación", "evita inferencia inconsistente"),
    ("P5", "NutriScore faltante", "'Falta Dato' → nulo estructural (no imputado)", "7 productos / 3,576 eventos"),
    ("P6", "Deduplicación", "clave event_id", "0 duplicados"),
]
top = Inches(1.7)
col_w = [Inches(0.9), Inches(2.6), Inches(5.3), Inches(3.3)]
left = Inches(0.5)
for j, (a, b, c, d) in enumerate([("Regla", "Nombre", "Qué hace", "Efecto medido")] + rows):
    x = left
    y = top + Inches(j * 0.62)
    header_row = j == 0
    for k, text in enumerate([a, b, c, d]):
        box = s.shapes.add_textbox(x, y, col_w[k], Inches(0.58))
        tfx = box.text_frame
        tfx.word_wrap = True
        px = tfx.paragraphs[0]
        px.text = text
        px.font.size = Pt(12.5 if not header_row else 12.5)
        px.font.bold = header_row or k == 0
        px.font.color.rgb = WHITE if header_row else (NAVY if k == 0 else INK)
        px.font.name = "Segoe UI"
        x += col_w[k]
    if header_row:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, sum(col_w, Emu(0)), Inches(0.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
        bar.shadow.inherit = False
        s.shapes._spTree.remove(bar._element)
        s.shapes._spTree.insert(2, bar._element)
add_footer(s, "docs/bitacora_entregas.md — Entrega 2 · verificado por reimplementación independiente (notebook 02 vs. preprocessing.py)")

# ---------------------------------------------------------------- Slide 7: Modeling
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Modelado predictivo", "Random Forest vs. Regresión Logística, con control de fuga de datos", RED)
if add_image_safe(s, FIG / "comparacion_roc.png", Inches(7.0), Inches(1.7), width=Inches(5.7)):
    pass
add_bullets(s, [
    ("Problema:", "clasificación binaria de eventos OUT → is_waste (Waste/Forced_Waste) vs. Consumption."),
    ("class_weight='balanced':", "el target está desbalanceado (~65% consumo / ~35% pérdida)."),
    ("Ganador:", "Random Forest — superior en F1, Recall y ROC-AUC."),
    ("Decisión técnica clave:", "dias_para_vencer resultó ser fuga de datos — el simulador la usaba para asignar las etiquetas (F1=1.0 trivial si se incluye). Se excluyó del entrenamiento; se conservó en el dataset para uso descriptivo."),
], top=1.7, width=6.2, size=15)
add_footer(s, "docs/tabla_comparativa_modelos.md · notebooks/03_modelo_metricas.ipynb")

# ---------------------------------------------------------------- Slide 8: PCA/t-SNE
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Componente avanzado", "PCA + t-SNE: el hallazgo que valida todo el dashboard", AMBER)
if add_image_safe(s, FIG / "tsne_scatter_2d.png", Inches(7.1), Inches(1.55), width=Inches(5.6)):
    pass
add_bullets(s, [
    ("Por qué:", "61 variables por evento; PCA (lineal, mide varianza retenida) + t-SNE (no lineal, agrupamientos locales) — complementarias, no intercambiables."),
    ("Resultado:", "29/61 componentes retienen 90% de varianza. Los productos se agrupan casi perfectamente por categoría."),
    ("Hallazgo de negocio:", "location es función determinista de category_name (verificado con crosstab) — cada categoría vive en una sola ubicación."),
    ("Impacto:", "el dashboard organiza todo por ubicación porque es, estructuralmente, la variable que más separa los datos — no una elección arbitraria."),
], top=1.6, width=6.4, size=14.5)
add_footer(s, "docs/componente_avanzado.md — justificación, mecánica, aplicación e interpretación completas")

# ---------------------------------------------------------------- Slide 9: Dashboard KPIs
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Dashboard", "KPIs comparativos — cada número, con su referencia", NAVY)
if not add_image_safe(s, SHOTS / "kpi.png", Inches(0.6), Inches(1.6), width=Inches(12.1)):
    stat_card(s, Inches(0.5), Inches(1.7), Inches(3.0), Inches(1.7), "Pérdida total (90 días)", "S/ 38,830", "17.5% del gasto en compras", RED)
    stat_card(s, Inches(3.65), Inches(1.7), Inches(3.0), Inches(1.7), "Tasa desperdicio Refrig.", "37.5%", "9.1x más que Estante", AMBER)
    stat_card(s, Inches(6.8), Inches(1.7), Inches(3.0), Inches(1.7), "Retraso al vencimiento", "-2.7 d", "sale del inventario ya vencido", RGBColor(0xF1, 0xC4, 0x0F))
    stat_card(s, Inches(9.95), Inches(1.7), Inches(3.0), Inches(1.7), "Saludables perdidos", "S/ 24,000", "45x más que en Despensa", GREEN)
add_bullets(s, [
    "17.5% del gasto total en compras terminó perdido — casi idéntico al 17% que estima el PNUMA/ONU a nivel mundial (Food Waste Index 2024).",
    "El Refrigerador concentra 97% de la pérdida económica.",
], top=4.1, size=16)
add_footer(s, "tableau/Sem12_Dashboard_SmartKitchen.html — cada KPI se calcula en vivo desde los datos, no está hardcodeado")

# ---------------------------------------------------------------- Slide 10: Dashboard charts
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Dashboard", "Evolución temporal y cruce con calidad nutricional", NAVY)
y_img = Inches(1.6)
placed = add_image_safe(s, SHOTS / "longitudinal.png", Inches(0.5), y_img, width=Inches(12.3))
if not placed:
    add_bullets(s, ["(Ver tableau/Sem12_Dashboard_SmartKitchen.html — sección 'Cuándo se pierde')"], top=2.5)
add_footer(s, "Leyendas al costado, títulos descriptivos, insight clave debajo de cada título")

# ---------------------------------------------------------------- Slide 11: Insights
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Insights", "Lo que dice el dashboard, en 5 líneas", GREEN)
add_bullets(s, [
    "El Refrigerador concentra 97% de la pérdida económica (S/ 37,568 de S/ 38,830).",
    "Tasa de desperdicio en Refrigerador: 37.5% — 9.1x más que Estante (4.1%), la ubicación más eficiente.",
    "S/ 24,000 en frutas/verduras saludables (NutriScore A-B) se pierden en el Refrigerador — 45x más que en Despensa. La pérdida no es comida chatarra.",
    "Los productos salen del inventario, en promedio, 2.7 días DESPUÉS de vencer.",
    "17.5% del gasto en compras perdido ≈ 17% del benchmark mundial (PNUMA/ONU 2024) — la simulación es una aproximación razonable a un problema real.",
], top=1.7, size=17)
add_footer(s, "docs/executive_summary.md §4")

# ---------------------------------------------------------------- Slide 12: Conclusions + recommendations
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Conclusiones y recomendaciones", "De los datos a la decisión", RED)
add_bullets(s, [
    ("Conclusión:", "responder \"¿dónde se pierde el dinero y la salud?\" con rigor exigió limpieza verificada por partida doble, modelado con control de fuga de datos, y una validación estructural (PCA) de la premisa del dashboard."),
    ("Recomendación 1:", "priorizar el Refrigerador en cualquier intervención — ahí vive el 97% de la pérdida."),
    ("Recomendación 2:", "adelantar la alerta de vencimiento a 3–5 días antes (rango_criticidad ya prototipado en Tableau), dado que hoy se actúa ~2.7 días tarde."),
    ("Recomendación 3:", "no tratar \"ubicación\" y \"categoría\" como palancas independientes con este dataset — son la misma variable."),
], top=1.7, size=16)
add_footer(s, "docs/defense_narrative.md §9-10")

# ---------------------------------------------------------------- Slide 13: Limitations + future work
s = prs.slides.add_slide(blank)
add_bg(s)
add_header(s, "Limitaciones y trabajo futuro", "Qué no dice este proyecto, y qué sigue", MUTED)
add_bullets(s, [
    ("Limitación:", "dataset sintético — simulación calibrada, no observaciones de hogares reales."),
    ("Limitación:", "location ≡ category_name en este dataset — no se puede aislar \"dónde se guarda\" de \"qué se guarda\"."),
    ("Limitación:", "catálogo sin productos NutriScore E; el bucket \"Riesgo (D-E)\" nunca se puebla con estos 50 productos."),
    ("Futuro:", "piloto con datos reales de hogares para desconfundir ubicación de categoría."),
    ("Futuro:", "empaquetar el .twbx final en Tableau Desktop y verificar filtros/parámetros en vivo."),
    ("Futuro:", "alerta operativa a 72 horas basada en rango_criticidad, no solo panel descriptivo."),
], top=1.7, size=15.5)
add_footer(s, "docs/QA_validation.md §10 · docs/defense_narrative.md §11")

# ---------------------------------------------------------------- Slide 14: Closing
s = prs.slides.add_slide(blank)
add_bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.8), Inches(11.5), Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True
p0 = tf.paragraphs[0]
p0.text = "Gracias"
p0.font.size = Pt(44)
p0.font.bold = True
p0.font.color.rgb = WHITE
p0.font.name = "Segoe UI"
p1 = tf.add_paragraph()
p1.text = "Preguntas — repositorio y documentación completa disponibles en el proyecto"
p1.font.size = Pt(16)
p1.font.color.rgb = RGBColor(0xB8, 0xC2, 0xE0)
p1.font.name = "Segoe UI"
p1.space_before = Pt(14)

prs.save(OUT)
print(f"Deck generado: {OUT}")
print(f"Slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
